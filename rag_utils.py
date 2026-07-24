import os
import faiss
import torch
import numpy as np
from dotenv import load_dotenv
from PIL import Image
from transformers import (
    CLIPModel,
    CLIPProcessor
)
from langchain_core.documents import Document
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_groq import ChatGroq
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
    UnstructuredWordDocumentLoader,
    UnstructuredPowerPointLoader
)
import fitz
from docx import Document as DocxDocument
from pptx import Presentation
from storage import Storage
import streamlit as st
load_dotenv()

GROQ_API_KEY = st.text_input(
                "Enter your GROQ API Key",
                type = "password"
                )

class RAG:
    def __init__(self):
        self.storage = Storage()
        self.embeddings = HuggingFaceEmbeddings(
            model_name="sentence-transformers/all-MiniLM-L6-v2"
        )
        self.vector_db = self.storage.load_vectorstore(self.embeddings)
        self.clip_model = CLIPModel.from_pretrained("openai/clip-vit-base-patch32")
        self.clip_processor = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch32")
        self.image_index = self.storage.load_image_index()
        self.image_metadata = self.storage.load_image_metadata()
        self.llm = ChatGroq(
            model="llama-3.1-8b-instant",
            temperature=0,
            api_key=GROQ_API_KEY
        )
        self.splitter = RecursiveCharacterTextSplitter(chunk_size=800,chunk_overlap=150)

    def extract_pdf_images(self, pdf_path):
        image_paths = []
        output_dir = os.path.join(
            self.storage.upload_folder,
            "images"
        )
        os.makedirs(output_dir, exist_ok=True)
        pdf = fitz.open(pdf_path)
        for page_num in range(len(pdf)):
            page = pdf.load_page(page_num)
            images = page.get_images(full=True)
            for img_num, img in enumerate(images):
                xref = img[0]
                base_image = pdf.extract_image(xref)
                image_bytes = base_image["image"]
                ext = base_image["ext"]
                image_path = os.path.join(
                    output_dir,
                    f"{os.path.splitext(os.path.basename(pdf_path))[0]}_page{page_num+1}_{img_num}.{ext}"
                )
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
                image_paths.append(image_path)
        pdf.close()
        return image_paths

    def extract_docx_images(self, docx_path):
        image_paths = []
        output_dir = os.path.join(
            self.storage.upload_folder,
            "images"
        )
        os.makedirs(output_dir, exist_ok=True)
        doc = DocxDocument(docx_path)
        image_no = 1
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image = rel.target_part.blob
                ext = rel.target_ref.split(".")[-1]
                image_path = os.path.join(
                    output_dir,
                    f"{os.path.splitext(os.path.basename(docx_path))[0]}_{image_no}.{ext}"
                )
                with open(image_path, "wb") as f:
                    f.write(image)
                image_paths.append(image_path)
                image_no += 1
        return image_paths

    def extract_ppt_images(self, ppt_path):
        image_paths = []
        output_dir = os.path.join(
            self.storage.upload_folder,
            "images"
        )
        os.makedirs(output_dir, exist_ok=True)
        prs = Presentation(ppt_path)
        image_no = 1
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:
                    image = shape.image
                    image_path = os.path.join(
                        output_dir,
                        f"{os.path.splitext(os.path.basename(ppt_path))[0]}_{image_no}.{image.ext}"
                    )
                    with open(image_path, "wb") as f:
                        f.write(image.blob)
                    image_paths.append(image_path)
                    image_no += 1
        return image_paths

    def load_document(self, path):
        if path.endswith(".pdf"):
            loader = PyPDFLoader(path)

        elif path.endswith(".txt"):
            loader = TextLoader(path)

        elif path.endswith(".docx"):
            loader = UnstructuredWordDocumentLoader(path)

        elif path.endswith(".pptx"):
            loader = UnstructuredPowerPointLoader(path)

        else:
            return []
        return loader.load()


    def save_table(self, table, filename, number):
        name = os.path.splitext(filename)[0]

        txt_path = os.path.join(
            self.storage.table_folder,
            f"{name}_table_{number}.txt"
        )

        with open(
            txt_path,
            "w",
            encoding="utf-8"
        ) as f:
            f.write(table)
        return table


    def get_image_embedding(self, image_path):
        image = Image.open(image_path).convert("RGB")
        inputs = self.clip_processor(
            images=image,
            return_tensors="pt"
        )

        with torch.no_grad():
            output = self.clip_model.get_image_features(**inputs)
            embedding = self.clip_model.get_image_features(**inputs)

        embedding = embedding[0].cpu().numpy()
        embedding = embedding / np.linalg.norm(embedding)

        return embedding.astype("float32")

    def get_text_embedding(self, text):
        inputs = self.clip_processor(
            text=[text],
            return_tensors="pt",
            padding=True
        )

        with torch.no_grad():
            embedding = self.clip_model.get_text_features(**inputs)
        embedding = embedding[0].cpu().numpy()
        embedding = embedding / np.linalg.norm(embedding)

        return embedding.astype("float32")

    def extract_images(self, file_path):
        if file_path.endswith(".pdf"):
            return self.extract_pdf_images(file_path)
        elif file_path.endswith(".docx"):
            return self.extract_docx_images(file_path)
        elif file_path.endswith(".pptx"):
            return self.extract_ppt_images(file_path)
        return []
 
    def process_file(self, file_path):
        docs = self.load_document(file_path)
        image_vectors = []
        images = self.extract_images(file_path)

        for image_path in images:
            vector = self.get_image_embedding(image_path)
            image_vectors.append(vector)
            self.image_metadata.append(
                {
                    "file": os.path.basename(file_path),
                    "path": image_path
                }
            )

        #for doc in docs:
            #print(doc.metadata)

        documents = []
        table_count = 1
        for doc in docs:
            metadata = doc.metadata.copy()
            metadata["file"] = os.path.basename(file_path)
            category = metadata.get("category", "Text")
            content = doc.page_content.strip()
            if not content:
                continue

            if category == "Table":
                content = self.save_table(
                    content,
                    metadata["file"],
                    table_count
                )
                table_count += 1
                doc_type = "table"

            else:
                doc_type = "text"
            chunks = self.splitter.split_text(content)

            for chunk in chunks:
                documents.append(
                    Document(
                        page_content=chunk,
                        metadata={
                                **metadata,
                                "type": doc_type
                            }
                        )
                    )
        return documents, image_vectors

    def index_files(self, file_paths):
        new_documents = []
        new_image_vectors = []

        for file_path in file_paths:
            docs, images = self.process_file(file_path)
            new_documents.extend(docs)
            new_image_vectors.extend(images)

        if new_documents:
            if self.vector_db is None:
                self.vector_db = FAISS.from_documents(new_documents,self.embeddings)

            else:
                self.vector_db.add_documents(new_documents)

            self.storage.save_vectorstore(self.vector_db)

        if new_image_vectors:
            vectors = np.array(
                new_image_vectors,
                dtype="float32"
            )

            if self.image_index is None:
                self.image_index = faiss.IndexFlatL2(vectors.shape[1])

            self.image_index.add(vectors)
            self.storage.save_image_index(self.image_index)
            self.storage.save_image_metadata(self.image_metadata)
 
    def search(self, query, k=5):
            text_results = []
            image_results = []

            if self.vector_db is not None:
                text_results = self.vector_db.similarity_search(
                    query,
                    k=k
                )

            if self.image_index is not None and len(self.image_metadata) > 0:
                query_vector = self.get_text_embedding(query).reshape(1, -1)
                _, indices = self.image_index.search(
                    query_vector,
                    min(k, len(self.image_metadata))
                )

                for idx in indices[0]:
                    if idx < len(self.image_metadata):
                        image_results.append(self.image_metadata[idx])

            return text_results, image_results

    def ask(self, query):
        history = self.storage.load_chat_history()
        query_embedding = self.embeddings.embed_query(query)
        best_answer = None
        best_score = -1

        for item in history:
            old_embedding = self.embeddings.embed_query(item["question"])
            score = np.dot(
                query_embedding,
                old_embedding
            ) / (
                np.linalg.norm(query_embedding)
                * np.linalg.norm(old_embedding)
            )

            if score > best_score:
                best_score = score
                best_answer = item["answer"]

        if best_score > 0.95:
            return {
                "answer": best_answer,
                "text": [],
                "images": [],
                "cached": True
            }

        text_results, image_results = self.search(query)
        context = ""

        for doc in text_results:
            if doc.metadata["type"] == "table":
                context += "Table:\n"
            else:
                context += "Text:\n"

            context += doc.page_content
            context += "\n\n"
        prompt = f"""
            You are an AI assistant.
            Answer ONLY from the given context.
            If the answer is not available, reply:
            I couldn't find that information in the uploaded documents.
            Context:
            {context}
            Question:
            {query}
            Answer:
            """

        response = self.llm.invoke(prompt)
        answer = response.content

        self.storage.save_chat(query,answer)

        return {
            "answer": answer,
            "text": text_results,
            "images": image_results,
            "cached": False
        }

    def delete_file(self, filename):

        if self.vector_db is not None:
            docs = []
            for _, doc in self.vector_db.docstore._dict.items():
                if doc.metadata.get("file") != filename:
                    docs.append(doc)

            if docs:
                self.vector_db = FAISS.from_documents(
                    docs,
                    self.embeddings
                )
                self.storage.save_vectorstore(self.vector_db)

            else:
                self.vector_db = None

        if self.image_metadata:
            vectors = []
            metadata = []

            for item in self.image_metadata:
                if item["file"] != filename:
                    metadata.append(item)
                    vectors.append(self.get_image_embedding(item["path"]))

            self.image_metadata = metadata
            self.storage.save_image_metadata(metadata)

            if vectors:
                vectors = np.array(vectors,dtype="float32")
                self.image_index = faiss.IndexFlatL2(vectors.shape[1])
                self.image_index.add(vectors)
                self.storage.save_image_index(self.image_index)

            else:
                self.image_index = None

        indexed = self.storage.load_indexed_files()

        if filename in indexed:
            del indexed[filename]
            self.storage.save_indexed_files(indexed)

        self.storage.delete_uploaded_file(filename)

    def clear_database(self):
        self.storage.clear_database()
        self.vector_db = None
        self.image_index = None
        self.image_metadata = []
 